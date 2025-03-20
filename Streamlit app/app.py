import streamlit as st
from PIL import Image
import pandas as pd
from pptx import Presentation
import PyPDF2
import docx
import logging

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Configure the Streamlit page
st.set_page_config(
    page_title="Document Redaction Tool",
    page_icon="🌎",
    layout="wide"
)

def main():
    # Sidebar
    with st.sidebar:
        # Instructions for using the app
        st.markdown("""
                    ## About
                    This document redaction tool allows you to upload files and redact personally identifiable information (PII).

                    ## Supported File Types:
                    - Text Files (.txt)
                    - PDF Files (.pdf)
                    - Word Documents (.docx)
                    - CSV Files (.csv)
                    - Image Files (.jpg, .jpeg, .png)
                    - PowerPoint Presentations (.pptx)

                    ## Supported Languages:
                    English

                    ### How to Use:
                    1. Upload a file 
                    2. You can view or redact PII.

                    """
                    )

    # Set the title of the app
    st.title("Document Redaction Tool")

    # Create an upload button to upload a document
    uploaded_file = st.file_uploader("Please upload a file...", type=["txt", "pdf", "docx", "jpg", "jpeg", "png", "gif", "csv", "pptx"])

    # Check if a file has been uploaded 
    if uploaded_file is not None:
        # Display the file name
        st.write(f"Filename: {uploaded_file.name}")
        
        # Try to handle the uploaded file based on its type
        try:
            # If the uploaded file is a text file
            if uploaded_file.type == "text/plain":
                content = uploaded_file.read().decode("utf-8")
                st.text_area("File Content", content, height=300)
            
            # If the uploaded file is a PDF file
            elif uploaded_file.type == "application/pdf":
                pdf_reader = PyPDF2.PdfReader(uploaded_file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text()
                st.text_area("File Content", text, height=300)
            
            # If the uploaded file is a Word document
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc = docx.Document(uploaded_file)
                text = ""
                for para in doc.paragraphs:
                    text += para.text + "\n"
                st.text_area("File Content", text, height=300)
            
            # If the uploaded file is an image
            elif uploaded_file.type in ["image/jpeg", "image/png", "image/gif"]:
                image = Image.open(uploaded_file)
                st.image(image, caption="Uploaded Image", use_container_width=True)
            
            # If the uploaded file is a CSV file
            elif uploaded_file.name.endswith(".csv") or uploaded_file.type == "text/csv":
                try:
                    df = pd.read_csv(uploaded_file)
                    st.write(df)  # Display CSV content as a table
                except Exception as e:
                    logging.error(f"Error reading CSV file: {e}")
                    st.error(f"Error reading CSV file: {e}")
            
            # If the uploaded file is a PowerPoint presentation
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                prs = Presentation(uploaded_file)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                st.text_area("File Content", text, height=300)
            
            # If the file type is unsupported
            else:
                raise ValueError("Unsupported file type.")
        
        except Exception as e:
            logging.error(f"Error processing file {uploaded_file.name}: {e}")
            st.error(f"Error processing file: {e}")

if __name__ == "__main__":
    main()
